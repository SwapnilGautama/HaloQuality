# share_snapshot.py
from __future__ import annotations
import io
import importlib
from datetime import datetime
from typing import Dict, Any, List, Tuple

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import smtplib
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email import encoders

# ---------------- PPTX builder ----------------

def _add_title(slide, title: str, subtitle: str | None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9)).text_frame
    tb.text = title
    p = tb.paragraphs[0]; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = RGBColor(11, 61, 145)
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.5)).text_frame
    stamp = datetime.now().strftime("%d %b %Y %H:%M")
    sb.text = f"{subtitle or ''}  •  {stamp}".strip()
    sb.paragraphs[0].font.size = Pt(12); sb.paragraphs[0].font.color.rgb = RGBColor(90,90,90)

def _add_fig(slide, fig, caption: str, x: float, y: float, w: float) -> float:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=220, bbox_inches="tight")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), width=Inches(w))
    cap = slide.shapes.add_textbox(Inches(x), Inches(y + 2.9), Inches(w), Inches(0.3)).text_frame
    cap.text = caption
    cap.paragraphs[0].font.size = Pt(10); cap.paragraphs[0].font.color.rgb = RGBColor(90,90,90)
    return y + 3.2

def _add_table_text(slide, df: pd.DataFrame, caption: str, x: float, y: float, w: float) -> float:
    txt = df.head(10).to_string(index=False)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.6)).text_frame
    box.text = caption
    box.paragraphs[0].font.size = Pt(12); box.paragraphs[0].font.bold = True
    p = box.add_paragraph(); p.text = txt; p.font.size = Pt(10)
    return y + 1.8

def build_one_pager(title: str, subtitle: str | None,
                    figs: List[Tuple[str, "matplotlib.figure.Figure"]] | None = None,
                    tables: List[Tuple[str, pd.DataFrame]] | None = None) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title, subtitle)

    x, y, w = 0.5, 1.6, 9.2
    if figs:
        for cap, fig in figs:
            y = _add_fig(slide, fig, cap, x, y, w)
    if tables:
        for cap, df in tables:
            y = _add_table_text(slide, df, cap, x, y, w)

    out = io.BytesIO(); prs.save(out); out.seek(0)
    return out.read()

# -------------- Optional content provider per-question --------------
QUESTION_MODULE_PREFIXES = ("questions", "core.questions")

def get_snapshot_content(slug: str, store: Dict[str, Any], params: Dict[str, Any]) -> Dict[str, Any]:
    """
    If the question module defines build_snapshot(store, params) -> {"title":..., "subtitle":..., "figs":[(cap, fig)], "tables":[(cap, df)]}
    we use it. Otherwise we fall back to a minimal one-pager.
    """
    last_exc = None
    for prefix in QUESTION_MODULE_PREFIXES:
        try:
            mod = importlib.import_module(f"{prefix}.{slug}")
            if hasattr(mod, "build_snapshot"):
                return mod.build_snapshot(store, params)
        except Exception as e:
            last_exc = e
            continue

    # Fallback one-pager (generic)
    title = f"Halo Quality — {slug.replace('_',' ').title()} Snapshot"
    subtitle = "Auto summary"
    head = []
    # sniff common keys in store to show something useful
    if "fpa" in store:
        fpa_df = store["fpa"]
        if isinstance(fpa_df, pd.DataFrame) and not fpa_df.empty:
            head.append(("FPA (preview)", fpa_df.head(8)))
    if "complaints" in store:
        comp_df = store["complaints"]
        if isinstance(comp_df, pd.DataFrame) and not comp_df.empty:
            head.append(("Complaints (preview)", comp_df.head(8)))
    if not head:
        head.append(("Info", pd.DataFrame({"note": ["No custom snapshot function; please add build_snapshot()."]})))
    return {"title": title, "subtitle": subtitle, "figs": [], "tables": head}

# -------------------- Outlook mailer --------------------

def send_via_outlook(to_addrs: list[str], subject: str, html_body: str,
                     attachment_name: str, attachment_bytes: bytes,
                     smtp_host: str, smtp_port: int, username: str, password: str, sender: str):
    msg = MIMEMultipart()
    msg["From"] = sender
    msg["To"] = ", ".join(to_addrs)
    msg["Subject"] = subject

    msg.attach(MIMEText(html_body, "html"))

    part = MIMEBase("application",
                    "vnd.openxmlformats-officedocument.presentationml.presentation")
    part.set_payload(attachment_bytes)
    encoders.encode_base64(part)
    part.add_header("Content-Disposition", f'attachment; filename="{attachment_name}"')
    msg.attach(part)

    with smtplib.SMTP(smtp_host, smtp_port) as server:
        server.starttls()
        server.login(username, password)
        server.sendmail(sender, to_addrs, msg.as_string())
