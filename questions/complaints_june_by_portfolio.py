# questions/complaints_june_by_portfolio.py
from __future__ import annotations

import os
import re
from io import BytesIO
from typing import Dict, List, Tuple, Optional

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt

# -------------------------------
# Optional integrations
# -------------------------------

_OPENAI_READY = False
try:
    import openai  # type: ignore
    if os.getenv("OPENAI_API_KEY"):
        openai.api_key = os.getenv("OPENAI_API_KEY")
        _OPENAI_READY = True
except Exception:
    _OPENAI_READY = False

_PPT_READY = False
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    _PPT_READY = True
except Exception:
    _PPT_READY = False


# -------------------------------
# Theme helpers
# -------------------------------

_DARK_BLUE = "#0b3d91"
_DARK_GREY = "#333333"
_SOFT_GREY = "#DDDDDD"
_PASTEL_LINE = "#88cde8"
_PARETO_LINE = "#6ab6e1"

def _header(title: str) -> None:
    st.markdown(
        f"<h3 style='color:{_DARK_BLUE};margin:0 0 .35rem 0; font-weight:700;'>{title}</h3>",
        unsafe_allow_html=True,
    )

def _hide_index(sty: "pd.io.formats.style.Styler") -> "pd.io.formats.style.Styler":
    try:
        return sty.hide(axis="index")
    except Exception:
        try:
            return sty.hide_index()
        except Exception:
            return sty

def _style_table(df: pd.DataFrame, formats: Optional[Dict[str, str]] = None) -> "pd.io.formats.style.Styler":
    sty = (
        df.style
        .set_table_styles(
            [
                {"selector": "th", "props": [("color", _DARK_BLUE), ("font-weight", "700")]},
                {"selector": "tbody td", "props": [("color", _DARK_GREY)]},
            ]
        )
        .set_properties(**{"color": _DARK_GREY})
    )
    sty = _hide_index(sty)
    sty = sty.format(formats) if formats else sty.format(precision=3)
    return sty


# -------------------------------
# Column detection & months
# -------------------------------

def _find_first_col(df: pd.DataFrame, candidates: List[str]) -> Optional[str]:
    cols = {c.lower(): c for c in df.columns}
    for c in candidates:
        if c.lower() in cols:
            return cols[c.lower()]
    return None

def _norm_month_from_series(s: pd.Series) -> pd.Series:
    if pd.api.types.is_period_dtype(s):
        return s.astype(str)
    dt = pd.to_datetime(s, errors="coerce", dayfirst=True, utc=False)
    return dt.dt.to_period("M").astype(str)

@st.cache_data(show_spinner=False)
def _build_month_column(df: pd.DataFrame, raw_col: str, assume_year: Optional[int] = None) -> pd.Series:
    s = df[raw_col]
    m = _norm_month_from_series(s)
    if m.notna().sum() >= max(1, int(0.1 * len(m))):
        return m
    if assume_year is not None:
        try:
            coerced = pd.to_datetime(s.astype(str) + f" {assume_year}", format="%B %Y", errors="coerce")
        except Exception:
            coerced = pd.to_datetime(s.astype(str) + f" {assume_year}", errors="coerce")
        return coerced.dt.to_period("M").astype(str)
    return pd.to_datetime(s, errors="coerce", dayfirst=True).dt.to_period("M").astype(str)

def _add_total_row(df: pd.DataFrame, sum_cols: List[str], label_col: str, label="Total") -> pd.DataFrame:
    total = {c: df[c].sum() if c in sum_cols else None for c in df.columns}
    total[label_col] = label
    out = pd.concat([pd.DataFrame([total]), df], ignore_index=True)
    if all(c in out.columns for c in ["cases", "complaints"]):
        with np.errstate(divide="ignore", invalid="ignore"):
            per = out["complaints"] * 1000 / out["cases"]
        out["per_1000"] = per.replace([np.inf, -np.inf], np.nan)
    return out

def _months_jan_to_aug_2025() -> List[str]:
    return [f"2025-{i:02d}" for i in range(1, 9)]

@st.cache_data(show_spinner=False)
def _latest_month_2025(cases: pd.DataFrame, comp: pd.DataFrame) -> Tuple[Optional[str], str]:
    c_date = _find_first_col(
        cases,
        ["Create Date (cases)", "Create Date", "Create date", "Start Date", "StartDate", "Created On", "CreateDt"],
    )
    k_date = _find_first_col(
        comp,
        ["Date Complaint Received - DD/MM/YY", "Date Complaint Received", "Complaint Date", "Received Date", "Month"],
    )
    if c_date is None and k_date is None:
        return None, "Latest"

    latest = None
    if c_date is not None:
        cm = pd.to_datetime(cases[c_date], errors="coerce", dayfirst=True).dt.to_period("M")
        cm = cm[cm.dt.year == 2025]
        if not cm.dropna().empty:
            latest = cm.max()

    if latest is None and k_date is not None:
        if k_date.lower() == "month":
            km = pd.to_datetime(comp[k_date].astype(str) + " 2025", errors="coerce").dt.to_period("M")
        else:
            km = pd.to_datetime(comp[k_date], errors="coerce", dayfirst=True).dt.to_period("M")
        km = km[km.dt.year == 2025]
        if not km.dropna().empty:
            latest = km.max()

    if latest is None:
        return None, "Latest"
    return str(latest), pd.Period(latest).to_timestamp().strftime("%b %Y")


# -------------------------------
# RCA labelling (rules + optional AI)
# -------------------------------

def _preclean(text: str) -> str:
    t = (text or "").lower()
    t = re.sub(r"[_/\\\-]+", " ", t)
    t = re.sub(r"[^a-z0-9 %]", " ", t)
    replacements = {
        r"\bbacs\b": " bank payment ",
        r"\bchaps\b": " bank payment ",
        r"\bdd\b": " direct debit ",
        r"\bsoa\b": " statement of account ",
        r"\btpa\b": " third party ",
        r"\bifa\b": " third party ",
        r"\bpoa\b": " power of attorney ",
        r"\baddr\b": " address ",
        r"\bni\b": " national insurance ",
        r"\bsort\b": " sort code ",
        r"\bacc\b": " account ",
        r"\brecalc(ulate|ulation|)\b": " recalculation ",
        r"\bcalc(ulate|ulation|)\b": " calculation ",
        r"\bresp\b": " response ",
        r"\bsig(n|)\b": " sign ",
        r"\bid\b": " identification ",
        r"\bkyc\b": " identification ",
        r"\bcof\b": " change of fund ",
    }
    for pat, repl in replacements.items():
        t = re.sub(pat, repl, t)
    t = re.sub(r"\s+", " ", t).strip()
    return t

def _classify(text: str, patterns: Dict[str, List[str]], default: str) -> str:
    t = _preclean(text)
    for label, pats in patterns.items():
        for p in pats:
            if re.search(p, t):
                return label
    return default

def _rca1_keyword(text: str) -> str:
    patterns = {
        "Delay": [
            r"\bdelay(ed|s|ing)?\b", r"\btimes? ?scale\b", r"\bsla\b", r"\bbacklog\b",
            r"\bqueue\b", r"\boverdue\b", r"\bawait(ing)?\b", r"\bchase(r|s|d|)?\b",
        ],
        "Procedure": [
            r"\bscheme rules?\b", r"\bprocedure\b", r"\bprocess\b",
            r"\btemplate\b", r"\bform\b", r"\bconsent form\b",
            r"\bdocument(ation)? (missing|not provided|not received)\b",
            r"\bevidence (missing|not provided|not received)\b",
        ],
        "Communication": [
            r"\bcommunicat(e|ion|ions)\b", r"\bemail\b", r"\bletter\b",
            r"\bphone|call\b", r"\bupdate\b", r"\bno (reply|response)\b",
            r"\bunclear\b", r"\bmis-?communicat",
        ],
        "System": [
            r"\bsystem\b", r"\bportal\b", r"\bworkflow\b", r"\bautomation\b",
            r"\bbug\b", r"\btechnical\b", r"\bit\b",
        ],
        "Incorrect/Incomplete information": [
            r"\bincorrect\b", r"\bwrong\b", r"\bincomplete\b", r"\berror\b",
            r"\bmis-?key(ed)?\b", r"\btypo\b", r"\bmisallocat(ed|ion)\b",
            r"\baddress\b", r"\bbank\b", r"\bdob\b", r"\bnational insurance\b", r"\bsort code\b", r"\baccount\b",
        ],
    }
    return _classify(text, patterns, default="Other")

def _rca2_keyword(text: str) -> str:
    patterns = {
        "Manual calculation": [r"\bmanual calculat", r"\bre-?calculat", r"\brecalculation\b"],
        "Data entry error": [
            r"\bmis-?key", r"\bkey(ing|ed) error\b", r"\btypo\b", r"\bwrong (amount|date|address|bank|ni|nino)\b",
            r"\bincorrect (amount|date|address|bank|ni|nino)\b", r"\bduplicate( case|)\b",
        ],
        "Documentation missing": [
            r"\bdocument(ation)? (missing|not provided|not received)\b",
            r"\bevidence (missing|not provided|not received)\b",
            r"\b(ids?|passport|birth|marriage|death) (cert|certificate)\b",
            r"\bproof (of )?(address|id|identity)\b",
        ],
        "Template/Form issue": [
            r"\bwrong form\b", r"\bincorrect form\b", r"\btemplate\b",
            r"\bform (not|in) (complete|completed|signed)\b", r"\bmissing signature\b", r"\bunsigned\b",
            r"\bcheckbox|tick box\b",
        ],
        "Waiting on member/TPA": [
            r"\bawait(ing)?\b", r"\bchase(r|s|d|ing)?\b", r"\bno (reply|response)\b",
            r"\bthird party\b", r"\bifa\b", r"\binsurer\b",
        ],
        "Bank/Payment issue": [
            r"\bbank\b", r"\b(bank )?payment\b", r"\brefund\b", r"\breturned\b",
            r"\bbounced\b", r"\bchaps\b", r"\bbacs\b", r"\bcheque\b", r"\baccount\b", r"\bsort code\b",
        ],
        "Overpayment": [r"\boverpayment\b"],
        "Address/Contact incorrect": [
            r"\baddress (wrong|incorrect|incomplete|change|updated?)\b",
            r"\bcontact (details|number|email) (wrong|incorrect|missing|change|updated?)\b",
        ],
        "Pension set up": [r"\b(pension|record) set ?up\b", r"\bsetup\b"],
        "Postal delay": [r"\bpost(al)?\b", r"\bmail\b"],
        "AVC": [r"\bavc\b"],
        "Case not created": [r"\bcase not created\b", r"\bnot raised\b"],
        "2nd review / QA": [r"\b(second|2nd) review\b", r"\bqa\b"],
        "Trustee": [r"\btrustee\b"],
        "Death benefits payout": [r"\bdeath benefit", r"\bbeneficiar(y|ies)\b"],
        "Drop in value / factor change": [r"\bfactor change\b", r"\bdrop in value\b", r"\bmarket\b", r"\bunit price\b"],
        "Scheme rules": [r"\bscheme rules?\b", r"\blegislation\b"],
        "Communication unclear": [r"\black of clarity\b", r"\bnot clear\b", r"\bconfus"],
    }
    return _classify(text, patterns, default="Other")

_RCA1_ALLOWED = [
    "Delay", "Procedure", "Communication", "System",
    "Incorrect/Incomplete information", "Other",
]
_RCA2_ALLOWED = [
    "Manual calculation", "Documentation missing", "Template/Form issue", "Data entry error",
    "Waiting on member/TPA", "Bank/Payment issue", "Address/Contact incorrect",
    "Pension set up", "Postal delay", "AVC", "Case not created", "2nd review / QA",
    "Trustee", "Death benefits payout", "Overpayment", "Drop in value / factor change",
    "Scheme rules", "Communication unclear", "Other",
]

RCA2_TO_RCA1_MAP = {
    "Manual calculation": "Delay",
    "Waiting on member/TPA": "Delay",
    "Postal delay": "Delay",
    "Case not created": "Delay",
    "2nd review / QA": "Delay",
    "Pension set up": "Delay",
    "Trustee": "Delay",
    "AVC": "Delay",
    "Overpayment": "Delay",
    "Death benefits payout": "Delay",
    "Bank/Payment issue": "Delay",

    "Documentation missing": "Procedure",
    "Template/Form issue": "Procedure",
    "Scheme rules": "Procedure",

    "Communication unclear": "Communication",

    "Data entry error": "Incorrect/Incomplete information",
    "Address/Contact incorrect": "Incorrect/Incomplete information",
    "Drop in value / factor change": "Incorrect/Incomplete information",
}

DELAY_RCA2 = {
    "Manual calculation", "Waiting on member/TPA", "Postal delay",
    "Case not created", "2nd review / QA", "Pension set up",
    "Trustee", "AVC", "Overpayment", "Death benefits payout",
    "Bank/Payment issue",
}
DELAY_EXTERNAL = {"Waiting on member/TPA", "Bank/Payment issue", "Postal delay", "Trustee", "AVC"}
DELAY_APTIA = {"Manual calculation", "Case not created", "2nd review / QA", "Pension set up", "Overpayment", "Data entry error"}

@st.cache_data(show_spinner=False)
def _ai_label_batch(texts: List[str]) -> List[Tuple[str, str]]:
    if not _OPENAI_READY:
        return [(_rca1_keyword(t), _rca2_keyword(t)) for t in texts]
    try:
        prompt = (
            "You are classifying complaint root causes. "
            "For each 'Brief Description – RCA done by admin', return a JSON array of objects "
            "with keys 'rca1' and 'rca2'.\n\n"
            f"RCA1 must be one of: {', '.join(_RCA1_ALLOWED)}.\n"
            f"RCA2 must be one of: {', '.join(_RCA2_ALLOWED)}.\n"
            "Prefer specific labels over 'Other'.\n\n"
            "Descriptions:\n" + "\n".join([f"- {t}" for t in texts])
        )
        resp = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Return only valid JSON."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.0,
        )
        content = resp["choices"][0]["message"]["content"]
        out: List[Tuple[str, str]] = []
        try:
            data = pd.read_json(content)
            if isinstance(data, pd.DataFrame):
                for _, row in data.iterrows():
                    r1 = row.get("rca1", "Other")
                    r2 = row.get("rca2", "Other")
                    r1 = r1 if r1 in _RCA1_ALLOWED else "Other"
                    r2 = r2 if r2 in _RCA2_ALLOWED else "Other"
                    out.append((r1, r2))
        except Exception:
            out = []
        if len(out) != len(texts):
            return [(_rca1_keyword(t), _rca2_keyword(t)) for t in texts]
        return out
    except Exception:
        return [(_rca1_keyword(t), _rca2_keyword(t)) for t in texts]


# -------------------------------
# Field detection
# -------------------------------

def _detect_cases_fields(cases: pd.DataFrame) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    id_col = _find_first_col(cases, ["Case ID", "CaseId", "ID"])
    port_col = _find_first_col(cases, ["Portfolio", "portfolio"])
    date_col = _find_first_col(
        cases, ["Create Date (cases)", "Create Date", "Create date", "Start Date", "StartDate", "Created On", "CreateDt"]
    )
    return id_col, port_col, date_col

def _detect_complaints_fields(comp: pd.DataFrame) -> Tuple[Optional[str], Optional[str], Optional[str], Optional[str]]:
    id_col = _find_first_col(comp, ["Original Process Affected Case ID", "Case ID", "Parent Case ID"])
    port_col = _find_first_col(comp, ["Portfolio", "portfolio"])
    date_col = _find_first_col(
        comp, [
            "Date Complaint Received - DD/MM/YY",
            "Date Complaint Received",
            "Complaint Date",
            "Received Date",
            "Month",
        ],
    )
    desc_col = _find_first_col(
        comp,
        ["Brief Description - RCA done by admin", "Brief Description", "RCA", "Admin Description", "Root Cause"],
    )
    return id_col, port_col, date_col, desc_col


# -------------------------------
# Core computations — cached
# -------------------------------

@st.cache_data(show_spinner=False)
def _portfolio_table_for_month(cases: pd.DataFrame, comp: pd.DataFrame, month_str: str) -> pd.DataFrame:
    id_c, port_c, date_c = _detect_cases_fields(cases)
    _, port_k, date_k, _ = _detect_complaints_fields(comp)
    if any(x is None for x in [port_c, date_c, port_k, date_k]) or not month_str:
        return pd.DataFrame(columns=["portfolio", "cases", "complaints", "per_1000"])

    cs = cases.copy(); kp = comp.copy()
    cs["_month"] = _build_month_column(cs, date_c)
    kp["_month"] = _build_month_column(kp, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)

    cases_m = cs.loc[cs["_month"] == month_str].groupby(port_c, dropna=False, as_index=False).size()
    cases_m.rename(columns={"size": "cases", port_c: "portfolio"}, inplace=True)

    comp_m = kp.loc[kp["_month"] == month_str].groupby(port_k, dropna=False, as_index=False).size()
    comp_m.rename(columns={"size": "complaints", port_k: "portfolio"}, inplace=True)

    out = pd.merge(cases_m, comp_m, how="left", on="portfolio")
    out["complaints"] = out["complaints"].fillna(0).astype(int)
    with np.errstate(divide="ignore", invalid="ignore"):
        out["per_1000"] = (out["complaints"] * 1000 / out["cases"]).replace([np.inf, -np.inf], np.nan)
    out["per_1000"] = out["per_1000"].round(1)

    out = out.sort_values(["complaints", "portfolio"], ascending=[False, True], kind="stable").reset_index(drop=True)
    out = _add_total_row(out, sum_cols=["cases", "complaints"], label_col="portfolio", label="Total")
    out["per_1000"] = out["per_1000"].round(1)
    return out

@st.cache_data(show_spinner=False)
def _mom_series(cases: pd.DataFrame, comp: pd.DataFrame) -> pd.DataFrame:
    _, _, date_c = _detect_cases_fields(cases)
    _, _, date_k, _ = _detect_complaints_fields(comp)
    if date_c is None or date_k is None:
        return pd.DataFrame(columns=["month", "per_1000"])

    cs = cases.copy(); kp = comp.copy()
    cs["_month"] = _build_month_column(cs, date_c)
    kp["_month"] = _build_month_column(kp, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)

    want = _months_jan_to_aug_2025()
    cs_m = cs.loc[cs["_month"].isin(want)].groupby("_month").size().reindex(want, fill_value=0)
    kp_m = kp.loc[kp["_month"].isin(want)].groupby("_month").size().reindex(want, fill_value=0)
    per_1000 = (kp_m * 1000 / cs_m.replace(0, np.nan)).fillna(0.0).round(1)
    pretty = [pd.Period(m).to_timestamp().strftime("%b-%y") for m in want]
    return pd.DataFrame({"month": pretty, "per_1000": per_1000.values})


# RCA tables for June
def _repair_rca1_from_rca2(rca1: List[str], rca2: List[str]) -> List[str]:
    return [RCA2_TO_RCA1_MAP.get(b, a) if a == "Other" else a for a, b in zip(rca1, rca2)]

@st.cache_data(show_spinner=False)
def _rca_tables_for_june(comp: pd.DataFrame, use_ai: bool) -> Tuple[pd.DataFrame, pd.DataFrame]:
    _, _, date_k, desc_col = _detect_complaints_fields(comp)
    if date_k is None or desc_col is None:
        return pd.DataFrame(columns=["RCA2", "count", "percent", "cum_percent"]), pd.DataFrame(columns=["RCA1", "count"])

    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)

    june = df.loc[df["_month"] == "2025-06", [desc_col]].fillna("")
    texts = june[desc_col].astype(str).tolist()

    if use_ai and _OPENAI_READY and len(texts) > 0:
        pairs: List[Tuple[str, str]] = []
        batch = 80
        for i in range(0, len(texts), batch):
            pairs.extend(_ai_label_batch(texts[i:i+batch]))
        r1_labels = [p[0] for p in pairs]
        r2_labels = [p[1] for p in pairs]
    else:
        r1_labels = [_rca1_keyword(t) for t in texts]
        r2_labels = [_rca2_keyword(t) for t in texts]

    r1_labels = _repair_rca1_from_rca2(r1_labels, r2_labels)

    r1 = pd.Series(r1_labels).value_counts(dropna=False).rename_axis("RCA1").reset_index(name="count")

    r2_counts = pd.Series(r2_labels).value_counts(dropna=False).rename_axis("RCA2").reset_index(name="count")
    r2_counts["order"] = np.where(r2_counts["RCA2"].eq("Other"), 1, 0)
    r2_counts = r2_counts.sort_values(["order", "count"], ascending=[True, False]).drop(columns="order").reset_index(drop=True)
    total = float(max(1, r2_counts["count"].sum()))
    r2_counts["percent"] = (r2_counts["count"] * 100 / total).round(1)
    r2_counts["cum_percent"] = r2_counts["percent"].cumsum().round(1)
    r2 = r2_counts.loc[r2_counts["cum_percent"] <= 80.0].reset_index(drop=True)

    return r2, r1

@st.cache_data(show_spinner=False)
def _rca2_table_by_portfolio_for_june(
    comp: pd.DataFrame,
    use_ai: bool,
    portfolios: Optional[List[str]] = None,
    rca1_keep: Optional[List[str]] = None,
) -> pd.DataFrame:
    _, port_k, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [port_k, date_k, desc_col]):
        return pd.DataFrame(columns=["Portfolio", "RCA2", "count"])

    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)
    df = df.loc[df["_month"] == "2025-06", [port_k, desc_col]].dropna(subset=[desc_col])
    if df.empty:
        return pd.DataFrame(columns=["Portfolio", "RCA2", "count"])

    if portfolios:
        df = df.loc[df[port_k].astype(str).isin([str(p) for p in portfolios])]
        if df.empty:
            return pd.DataFrame(columns=["Portfolio", "RCA2", "count"])

    texts = df[desc_col].astype(str).tolist()
    if use_ai and _OPENAI_READY and len(texts) > 0:
        pairs: List[Tuple[str, str]] = []
        batch = 80
        for i in range(0, len(texts), batch):
            pairs.extend(_ai_label_batch(texts[i:i+batch]))
        r1_labels = [p[0] for p in pairs]
        r2_labels = [p[1] for p in pairs]
    else:
        r1_labels = [_rca1_keyword(t) for t in texts]
        r2_labels = [_rca2_keyword(t) for t in texts]

    r1_labels = _repair_rca1_from_rca2(r1_labels, r2_labels)

    if rca1_keep:
        keep = set(rca1_keep)
        mask = [r in keep for r in r1_labels]
        if not any(mask):
            return pd.DataFrame(columns=["Portfolio", "RCA2", "count"])
        df = df.loc[mask].copy()
        r2_labels = [r2 for r2, m in zip(r2_labels, mask) if m]

    df["RCA2"] = r2_labels
    df["Portfolio"] = df[port_k].astype(str).fillna("")
    tab = df.groupby(["Portfolio", "RCA2"], dropna=False, as_index=False).size().rename(columns={"size": "count"})
    tab = tab.sort_values(["count", "Portfolio", "RCA2"], ascending=[False, True, True], kind="stable").reset_index(drop=True)
    return tab


# -------------------------------
# Plotting
# -------------------------------

def _mom_line_fig(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(6.0, 3.2))
    ax.plot(df["month"], df["per_1000"], marker="o", linewidth=2.5, color="#9ecae1")
    for x, y in zip(df["month"], df["per_1000"]):
        ax.text(x, y + 0.03, f"{y:.1f}", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)
    ax.set_title("Complaints per 1,000 — MoM (Jan–Aug ’25)", pad=8, color=_DARK_BLUE)
    ax.set_ylim(bottom=0)
    for sp in ["left", "right", "top"]:
        ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color(_SOFT_GREY); ax.spines["bottom"].set_linewidth(1.25)
    ax.get_yaxis().set_visible(False)
    ax.tick_params(axis="x", colors=_DARK_GREY)
    ax.set_xlabel(""); ax.set_ylabel(""); ax.grid(False)
    return fig

def _plot_mom_line(df: pd.DataFrame):
    st.pyplot(_mom_line_fig(df))

def _pareto_fig(df: pd.DataFrame):
    data = df.copy().sort_values("count", ascending=False).reset_index(drop=True)
    total = float(max(1, data["count"].sum()))
    data["percent"] = data["count"] * 100.0 / total
    data["cum_percent"] = data["percent"].cumsum()

    fig, ax = plt.subplots(figsize=(6.4, 3.8))
    bar_colors = ["#9ecae1", "#a1d99b", "#bdbdbd", "#fdd0a2", "#fdae6b", "#c7c7c7", "#bcbddc", "#ccebc5"]
    ax.bar(data["RCA1"], data["count"], color=bar_colors[: len(data)])
    for i, y in enumerate(data["count"].tolist()):
        ax.text(i, y + max(1, y * 0.02), f"{int(y)}", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)
    ax.set_title("RCA1 — June 2025 (Pareto)", pad=8, color=_DARK_BLUE)

    ax2 = ax.twinx()
    ax2.plot(data["RCA1"], data["cum_percent"], color=_PARETO_LINE, marker="o", linewidth=2.5)
    for i, y in enumerate(data["cum_percent"].tolist()):
        ax2.text(i, y + 1, f"{y:.0f}%", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)

    for sp in ["top", "right", "left"]:
        ax.spines[sp].set_visible(False)
    ax.grid(False); ax.get_yaxis().set_visible(False)
    ax.spines["bottom"].set_color(_SOFT_GREY); ax.spines["bottom"].set_linewidth(1.25)
    ax.tick_params(axis="x", colors=_DARK_GREY)
    plt.setp(ax.get_xticklabels(), rotation=90, ha="center")

    for sp in ["top", "right", "left", "bottom"]:
        ax2.spines[sp].set_visible(False)
    ax2.set_ylim(0, 100); ax2.set_ylabel(""); ax2.tick_params(axis="y", length=0); ax2.get_yaxis().set_visible(False)
    ax2.grid(False)
    return fig

def _plot_rca1_pareto(df: pd.DataFrame):
    st.pyplot(_pareto_fig(df))

def _fig_portfolio_trend(df: pd.DataFrame, title: str):
    fig, ax = plt.subplots(figsize=(6.8, 3.2))
    ax.plot(df["month"], df["per_1000"], color="#d95f02", linewidth=2.5, marker="o")
    for x, y in zip(df["month"], df["per_1000"]):
        ax.text(x, y + 0.15, f"{y:.1f}", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)
    idx = np.arange(len(df))
    z = np.polyfit(idx, df["per_1000"].astype(float), 1)
    ax.plot(df["month"], z[0]*idx + z[1], linestyle=":", linewidth=1.8, color="#7f7f7f")
    ax.set_title(title, color=_DARK_BLUE); ax.set_ylim(bottom=0)
    for sp in ["left", "right", "top"]:
        ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color(_SOFT_GREY); ax.spines["bottom"].set_linewidth(1.25)
    ax.get_yaxis().set_visible(False); ax.tick_params(axis="x", rotation=0, colors=_DARK_GREY)
    ax.set_xlabel(""); ax.set_ylabel(""); ax.grid(False)
    return fig

def _fig_reason_trend(df: pd.DataFrame):
    if df.empty:
        return None
    fig, ax = plt.subplots(figsize=(6.6, 3.2))
    x = np.arange(len(df["RCA1"])); width = 0.26
    colors = ["#74c476", "#a1d99b", "#9ecae1"]
    labels = [c for c in df.columns if c != "RCA1"]
    for i, col in enumerate(labels):
        ax.bar(x + (i-1)*width, df[col].values, width=width, label=col, color=colors[i % len(colors)])
    ax.set_xticks(x); ax.set_xticklabels(df["RCA1"], rotation=0, color=_DARK_GREY)
    ax.set_ylim(0, 100); ax.legend(frameon=False, loc="upper right")
    ax.set_title("Reason Trend (Jan–Aug ’25) — % split", color=_DARK_BLUE)
    for sp in ["left", "right", "top"]:
        ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color(_SOFT_GREY); ax.spines["bottom"].set_linewidth(1.25)
    ax.get_yaxis().set_visible(False); ax.set_ylabel(""); ax.set_xlabel("")
    return fig

def _fig_delay_split(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(6.6, 3.0))
    x = np.arange(len(df["month"])); width = 0.35
    ax.bar(x - width/2, df["External delay %"], width=width, color="#74c476", label="External Delay")
    ax.bar(x + width/2, df["Aptia delay %"], width=width, color="#9ecae1", label="Aptia Delay")
    for i, (e, a) in enumerate(zip(df["External delay %"], df["Aptia delay %"])):
        ax.text(i - width/2, e + 1, f"{e:.0f}%", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)
        ax.text(i + width/2, a + 1, f"{a:.0f}%", ha="center", va="bottom", fontsize=9, color=_DARK_GREY)
    ax.set_xticks(x); ax.set_xticklabels(df["month"], color=_DARK_GREY)
    ax.set_ylim(0, 100); ax.legend(frameon=False, loc="upper right")
    ax.set_title("Delay split — External vs Aptia (Jan–Aug ’25)", color=_DARK_BLUE)
    for sp in ["left", "right", "top"]:
        ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color(_SOFT_GREY); ax.spines["bottom"].set_linewidth(1.25)
    ax.get_yaxis().set_visible(False); ax.set_xlabel(""); ax.set_ylabel(""); ax.grid(False)
    return fig


# -------------------------------
# Portfolio-wise helpers
# -------------------------------

@st.cache_data(show_spinner=False)
def _portfolio_list(cases: pd.DataFrame, comp: pd.DataFrame) -> List[str]:
    _, port_c, _ = _detect_cases_fields(cases)
    _, port_k, _, _ = _detect_complaints_fields(comp)
    ports: set = set()
    if port_c and port_c in cases.columns:
        ports |= set(cases[port_c].dropna().astype(str).unique().tolist())
    if port_k and port_k in comp.columns:
        ports |= set(comp[port_k].dropna().astype(str).unique().tolist())
    desired = ["Chichester", "London", "Northwest", "Scotland"]
    return [p for p in desired if p in ports]

@st.cache_data(show_spinner=False)
def _portfolio_mom_series(cases: pd.DataFrame, comp: pd.DataFrame, portfolio: str) -> pd.DataFrame:
    _, port_c, date_c = _detect_cases_fields(cases)
    _, port_k, date_k, _ = _detect_complaints_fields(comp)
    if any(x is None for x in [port_c, date_c, port_k, date_k]):
        return pd.DataFrame(columns=["month", "per_1000"])

    cs = cases.copy(); kp = comp.copy()
    cs["_month"] = _build_month_column(cs, date_c)
    kp["_month"] = _build_month_column(kp, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)

    months = _months_jan_to_aug_2025()
    csz = cs.loc[(cs[port_c] == portfolio) & (cs["_month"].isin(months))].groupby("_month").size().reindex(months, fill_value=0)
    kpz = kp.loc[(kp[port_k] == portfolio) & (kp["_month"].isin(months))].groupby("_month").size().reindex(months, fill_value=0)
    per = (kpz * 1000 / csz.replace(0, np.nan)).fillna(0.0).values
    pretty = [pd.Period(m).to_timestamp().strftime("%b-%y") for m in months]
    return pd.DataFrame({"month": pretty, "per_1000": np.round(per, 1)})

@st.cache_data(show_spinner=False)
def _reason_trend_df(comp: pd.DataFrame, portfolio: str, use_ai: bool) -> pd.DataFrame:
    _, port_k, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [port_k, date_k, desc_col]):
        return pd.DataFrame()

    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)
    months = _months_jan_to_aug_2025()
    lab_months = [pd.Period(m).to_timestamp().strftime("%b’%y") for m in months]

    out: Dict[str, List[int]] = {"RCA1": ["Delay","Procedure","Communication","System","Incorrect/Incomplete information"]}
    for m, label in zip(months, lab_months):
        sub = df.loc[(df[port_k] == portfolio) & (df["_month"] == m), [desc_col]]
        if sub.empty:
            out[label] = [0,0,0,0,0]; continue
        texts = sub[desc_col].astype(str)
        r1 = [_rca1_keyword(t) for t in texts] if not (_OPENAI_READY and use_ai) else [p[0] for p in _ai_label_batch(list(texts))]
        s = pd.Series(r1).value_counts()
        total = max(1, s.sum())
        vals = [round(100*s.get(k, 0)/total) for k in ["Delay","Procedure","Communication","System","Incorrect/Incomplete information"]]
        out[label] = vals
    return pd.DataFrame(out)

@st.cache_data(show_spinner=False)
def _delay_split_df(comp: pd.DataFrame, portfolio: str, use_ai: bool) -> pd.DataFrame:
    _, port_k, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [port_k, date_k, desc_col]):
        return pd.DataFrame()

    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)

    months = _months_jan_to_aug_2025()
    labels = [pd.Period(m).to_timestamp().strftime("%b’%y") for m in months]

    rows: List[Tuple[str, float, float]] = []
    for m, lab in zip(months, labels):
        sub = df.loc[(df[port_k] == portfolio) & (df["_month"] == m), [desc_col]]
        if sub.empty:
            rows.append((lab, 0.0, 0.0)); continue
        texts = sub[desc_col].astype(str).tolist()
        if _OPENAI_READY and use_ai:
            pairs = _ai_label_batch(texts)
            r1 = [p[0] for p in pairs]; r2 = [p[1] for p in pairs]
        else:
            r1 = [_rca1_keyword(t) for t in texts]; r2 = [_rca2_keyword(t) for t in texts]
        delay_mask = [a == "Delay" or b in DELAY_RCA2 for a, b in zip(r1, r2)]
        if not any(delay_mask):
            rows.append((lab, 0.0, 0.0)); continue
        r2_delay = [r2[i] for i, keep in enumerate(delay_mask) if keep]
        total = len(r2_delay)
        ext = sum(1 for v in r2_delay if v in DELAY_EXTERNAL)
        apt = sum(1 for v in r2_delay if v in DELAY_APTIA)
        rows.append((lab, round(100*ext/total, 0) if total else 0.0, round(100*apt/total, 0) if total else 0.0))
    return pd.DataFrame(rows, columns=["month", "External delay %", "Aptia delay %"])


# -------------------------------
# PPT helper
# -------------------------------

def _add_df_table_to_slide(slide, df: pd.DataFrame, left_in: float, top_in: float, width_in: float):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.8 + 0.3*rows)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j); cell.text = str(col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(11, 61, 145)
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, col in enumerate(df.columns):
            val = row[col]
            if isinstance(val, float):
                if col == "per_1000":
                    text = f"{val:.1f}"
                else:
                    s = f"{val:.3f}"; text = s.rstrip('0').rstrip('.') if '.' in s else f"{val:.0f}"
            else:
                text = str(val)
            cell = table.cell(i, j); cell.text = text
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(11); p.font.color.rgb = RGBColor(51, 51, 51)
            p.alignment = PP_ALIGN.LEFT
    return table

def _build_ppt(table_df: pd.DataFrame, mom_df: pd.DataFrame, rca1_df: pd.DataFrame, rca2_df: pd.DataFrame) -> bytes:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    title = s1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)).text_frame
    title.text = "Complaint analysis — Jun 2025"; title.paragraphs[0].font.color.rgb = RGBColor(11, 61, 145); title.paragraphs[0].font.size = Pt(28)

    _add_df_table_to_slide(s1, table_df, left_in=0.5, top_in=1.2, width_in=5.0)
    fig_mom = _mom_line_fig(mom_df); buf_mom = BytesIO(); fig_mom.savefig(buf_mom, format="png", dpi=220, bbox_inches="tight"); plt.close(fig_mom)
    s1.shapes.add_picture(buf_mom, Inches(6.0), Inches(1.0), width=Inches(4.5))

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    t2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)).text_frame
    t2.text = "June reasons — RCA"; t2.paragraphs[0].font.color.rgb = RGBColor(11, 61, 145); t2.paragraphs[0].font.size = Pt(24)

    fig_p = _pareto_fig(rca1_df); buf_p = BytesIO(); fig_p.savefig(buf_p, format="png", dpi=220, bbox_inches="tight"); plt.close(fig_p)
    s2.shapes.add_picture(buf_p, Inches(0.5), Inches(1.1), width=Inches(5.6))
    _add_df_table_to_slide(s2, rca2_df, left_in=6.4, top_in=1.1, width_in=4.2)

    out = BytesIO(); prs.save(out); out.seek(0)
    return out.read()


# -------------------------------
# Insights helpers
# -------------------------------

@st.cache_data(show_spinner=False)
def _overall_year_rollup(cases: pd.DataFrame, comp: pd.DataFrame) -> pd.DataFrame:
    _, port_c, date_c = _detect_cases_fields(cases)
    _, port_k, date_k, _ = _detect_complaints_fields(comp)
    if any(x is None for x in [port_c, date_c, port_k, date_k]):
        return pd.DataFrame(columns=["Portfolio","Month","Cases","Complaints","Per1000"])
    months = _months_jan_to_aug_2025()
    cs = cases.copy(); cp = comp.copy()
    cs["_month"] = _build_month_column(cs, date_c)
    cp["_month"] = _build_month_column(cp, date_k, assume_year=2025 if date_k and date_k.lower()=="month" else None)
    cs = (cs[cs["_month"].isin(months)].groupby([port_c,"_month"]).size().rename("Cases").reset_index())
    cp = (cp[cp["_month"].isin(months)].groupby([port_k,"_month"]).size().rename("Complaints").reset_index())
    m = pd.merge(cs, cp, left_on=[port_c,"_month"], right_on=[port_k,"_month"], how="left")
    m.rename(columns={port_c:"Portfolio","_month":"Month"}, inplace=True)
    m["Complaints"] = m["Complaints"].fillna(0).astype(int)
    m["Per1000"] = (m["Complaints"] * 1000 / m["Cases"].replace(0,np.nan)).astype(float)
    return m

def _slope_pp_per_month(series: pd.Series) -> Optional[float]:
    y = pd.to_numeric(series, errors="coerce").dropna()
    if len(y) < 3: return None
    x = np.arange(len(y), dtype=float)
    try:
        b1, _ = np.polyfit(x, y.values, 1); return float(b1)
    except Exception:
        return None

@st.cache_data(show_spinner=False)
def _rca_year_counts(comp: pd.DataFrame, use_ai: bool) -> pd.DataFrame:
    _, _, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [date_k, desc_col]):
        return pd.DataFrame(columns=["Month","RCA1","RCA2","count"])
    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower()=="month" else None)
    df = df[df["_month"].isin(_months_jan_to_aug_2025())]
    if df.empty: return pd.DataFrame(columns=["Month","RCA1","RCA2","count"])
    if _OPENAI_READY and use_ai:
        pairs = _ai_label_batch(df[desc_col].astype(str).tolist())
        r1 = [p[0] for p in pairs]; r2 = [p[1] for p in pairs]
    else:
        r1 = [_rca1_keyword(t) for t in df[desc_col].astype(str)]
        r2 = [_rca2_keyword(t) for t in df[desc_col].astype(str)]
    out = pd.DataFrame({"Month": df["_month"].values, "RCA1": r1, "RCA2": r2})
    out["count"] = 1
    return out.groupby(["Month","RCA1","RCA2"], as_index=False)["count"].sum()

@st.cache_data(show_spinner=False)
def _delay_attribution_share(comp: pd.DataFrame, use_ai: bool) -> Tuple[float, float]:
    _, _, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [date_k, desc_col]): return (np.nan, np.nan)
    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower()=="month" else None)
    df = df[df["_month"].isin(_months_jan_to_aug_2025())]
    if df.empty: return (np.nan, np.nan)
    if _OPENAI_READY and use_ai:
        pairs = _ai_label_batch(df[desc_col].astype(str).tolist())
        r1 = [p[0] for p in pairs]; r2 = [p[1] for p in pairs]
    else:
        r1 = [_rca1_keyword(t) for t in df[desc_col].astype(str)]
        r2 = [_rca2_keyword(t) for t in df[desc_col].astype(str)]
    mask = [a == "Delay" or b in DELAY_RCA2 for a,b in zip(r1,r2)]
    if not any(mask): return (0.0, 0.0)
    r2_delay = [r2[i] for i, keep in enumerate(mask) if keep]
    total = len(r2_delay); ext = sum(1 for v in r2_delay if v in DELAY_EXTERNAL); apt = sum(1 for v in r2_delay if v in DELAY_APTIA)
    return round(100*ext/total,1), round(100*apt/total,1)

@st.cache_data(show_spinner=False)
def _combo_predictions(comp: pd.DataFrame, cases: pd.DataFrame, use_ai: bool) -> pd.DataFrame:
    roll = _overall_year_rollup(cases, comp)
    bench = (roll.groupby("Portfolio")["Per1000"].mean(skipna=True).rename("Per1000_portfolio").reset_index())
    _, port_k, date_k, desc_col = _detect_complaints_fields(comp)
    if any(x is None for x in [port_k, date_k, desc_col]) or roll.empty:
        return pd.DataFrame(columns=["Portfolio","RCA2","Complaints_2025","Per1000_portfolio"])
    df = comp.copy()
    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower()=="month" else None)
    df = df[df["_month"].isin(_months_jan_to_aug_2025())]
    if df.empty: return pd.DataFrame(columns=["Portfolio","RCA2","Complaints_2025","Per1000_portfolio"])
    r2 = [_rca2_keyword(t) for t in df[desc_col].astype(str)] if not (_OPENAI_READY and use_ai) else [p[1] for p in _ai_label_batch(df[desc_col].astype(str).tolist())]
    tab = (pd.DataFrame({"Portfolio": df[port_k].astype(str).values, "RCA2": r2})
           .assign(count=1)
           .groupby(["Portfolio","RCA2"], as_index=False)["count"].sum()
           .rename(columns={"count":"Complaints_2025"}))
    out = tab.merge(bench, on="Portfolio", how="left").sort_values(
        ["Complaints_2025","Per1000_portfolio","Portfolio","RCA2"],
        ascending=[False, False, True, True],
        kind="stable"
    )
    return out.head(5)


# -------------------------------
# Streamlit entrypoint
# -------------------------------

def run(store: Dict, params: Dict, user_text: str = "") -> Tuple[str, pd.DataFrame]:
    cases: pd.DataFrame = store.get("cases", pd.DataFrame()).copy()
    comp: pd.DataFrame = store.get("complaints", pd.DataFrame()).copy()
    use_ai = bool(os.getenv("OPENAI_API_KEY"))

    latest_month_str, latest_label = _latest_month_2025(cases, comp)

    portfolios = _portfolio_list(cases, comp)
    tabs = st.tabs(["Insights", "Overall"] + portfolios)

    # --------- INSIGHTS ---------
    with tabs[0]:
        st.markdown(f"<h2 style='color:{_DARK_BLUE};margin:.3rem 0 1rem 0;'>What’s happening and why</h2>", unsafe_allow_html=True)

        overall_roll = _overall_year_rollup(cases, comp)
        if overall_roll.empty:
            st.info("Not enough data to compute insights.")
        else:
            nat = overall_roll.groupby("Month")[["Cases","Complaints"]].sum().reset_index().sort_values("Month")
            nat["Per1000"] = (nat["Complaints"] * 1000 / nat["Cases"].replace(0,np.nan)).astype(float)
            slope = _slope_pp_per_month(nat["Per1000"])
            slope_txt = "n/a" if slope is None else f"{slope:+.2f} per 1,000 per month"
            latest_month = nat["Month"].max(); latest_p1000 = float(nat.loc[nat["Month"]==latest_month,"Per1000"].values[0])

            bench = (overall_roll.groupby("Portfolio")["Per1000"].mean(skipna=True).rename("Avg Per1000"))
            latest = (overall_roll.sort_values("Month").groupby("Portfolio").tail(1)
                      .set_index("Portfolio")["Per1000"].rename("Latest Per1000"))
            hot = pd.concat([bench, latest], axis=1).reset_index()
            hot["Δ latest vs avg"] = (hot["Latest Per1000"] - hot["Avg Per1000"]).round(1)
            hotspots = hot.sort_values(["Avg Per1000","Latest Per1000"], ascending=False, kind="stable").head(5)

            rca_year = _rca_year_counts(comp, use_ai=use_ai)
            story_rca = ""
            if not rca_year.empty:
                rca_m = (rca_year.groupby(["Month","RCA1"])["count"].sum()
                         .groupby(level=0).apply(lambda s: s/s.sum()*100).reset_index(name="pct"))
                early = rca_m[rca_m["Month"].isin([f"2025-{i:02d}" for i in range(1,5)])]
                late  = rca_m[rca_m["Month"].isin([f"2025-{i:02d}" for i in range(5,9)])]
                r_early = early.groupby("RCA1")["pct"].mean(); r_late = late.groupby("RCA1")["pct"].mean()
                diff = (r_late - r_early).dropna().sort_values(ascending=False)
                rising = ", ".join([f"{k} (+{v:.1f}pp)" for k,v in diff.head(3).items()]) if not diff.empty else "—"
                falling = ", ".join([f"{k} ({v:.1f}pp)" for k,v in diff.tail(3).items()]) if len(diff)>=3 else "—"
                story_rca = f"**Rising reasons**: {rising}  \n**Falling reasons**: {falling}"

            ext_pct, apt_pct = _delay_attribution_share(comp, use_ai=use_ai)

            top_rca2_df = pd.DataFrame(columns=["RCA2","Share %"])
            if not rca_year.empty:
                top_rca2 = (rca_year.groupby("RCA2")["count"].sum().sort_values(ascending=False).head(6))
                if not top_rca2.empty and top_rca2.sum() > 0:
                    top_rca2_df = ((top_rca2 / top_rca2.sum() * 100).round(1).rename("Share %").reset_index())

            combos = _combo_predictions(comp, cases, use_ai=use_ai)

            nat_avg = float(bench.mean(skipna=True)) if not bench.empty else float("nan")
            bench_flag = (bench.reset_index()
                               .rename(columns={"Per1000":"Avg Per1000"})
                               .assign(Benchmark=lambda d: d["Avg Per1000"].apply(lambda v: "Above" if pd.notna(nat_avg) and v>nat_avg else "Below")))

            colA, colB = st.columns((1.15, 1.0))
            with colA:
                st.markdown("#### At a glance")
                trend_txt = ("flat" if (slope is None or abs(slope) < 0.1) else ("upward" if slope>0 else "downward"))
                st.markdown(
                    f"- **Latest month:** `{latest_month}` • **National complaints/1,000:** **{latest_p1000:.1f}**\n"
                    f"- **Trend:** {trend_txt} *(slope {slope_txt})*\n"
                    f"- **Delay attribution (Jan–Aug ’25):** External **{ext_pct:.1f}%**, Aptia **{apt_pct:.1f}%**\n"
                    f"- **RCA movement:** {story_rca if story_rca else '—'}"
                )

                st.markdown("#### Complaint intensity hotspots — top 5 portfolios")
                st.dataframe(
                    _style_table(
                        hotspots[["Portfolio","Avg Per1000","Latest Per1000","Δ latest vs avg"]].round(1),
                        formats={"Avg Per1000":"{:.1f}","Latest Per1000":"{:.1f}","Δ latest vs avg":"{:+.1f}"}
                    ),
                    use_container_width=True,
                )

                st.markdown("#### Predicted high-risk combinations — Portfolio × RCA2 (by volume)")
                if combos.empty:
                    st.info("No combinations to highlight for Jan–Aug ’25.")
                else:
                    st.dataframe(
                        _style_table(
                            combos.rename(columns={"Complaints_2025":"Complaints (’25)","Per1000_portfolio":"Per1000 (portfolio avg)"}),
                            formats={"Complaints (’25)":"{:,.0f}","Per1000 (portfolio avg)":"{:.1f}"}
                        ),
                        use_container_width=True,
                    )

            with colB:
                st.markdown("#### Key issue themes — RCA2 (Jan–Aug ’25)")
                st.dataframe(
                    _style_table(top_rca2_df, formats={"Share %":"{:.1f}%"}),
                    use_container_width=True,
                )

                st.markdown("#### Above/below benchmark portfolios")
                st.dataframe(
                    _style_table(bench_flag.rename(columns={"Avg Per1000":"Avg Per1000 (’25)"}), formats={"Avg Per1000 (’25)":"{:.1f}"}),
                    use_container_width=True,
                )

            st.caption("Notes: Per-1,000 rates use cases as denominator. Trends use Jan–Aug ’25. RCA labels use rules with optional AI assist; themes are directional rather than exact.")

    # --------- OVERALL ---------
    with tabs[1]:
        table = _portfolio_table_for_month(cases, comp, latest_month_str or "2025-06")
        mom = _mom_series(cases, comp)

        c1, c2 = st.columns((1.2, 1.0), gap="large")
        with c1:
            _header(f"Complaint analysis — {latest_label} (by portfolio)")
            if table.empty: st.info("No rows returned for the current filters.")
            else:
                st.dataframe(_style_table(table, formats={"per_1000": "{:.1f}", "cases": "{:,.0f}", "complaints": "{:,.0f}"}), use_container_width=True)
        with c2:
            if not mom.empty:
                _plot_mom_line(mom)

        # Local filters + RCA visuals
        _, port_k, _, _ = _detect_complaints_fields(comp)
        rca1_options = _RCA1_ALLOWED
        ports_options = portfolios if portfolios else []

        f1, f2 = st.columns((1.0, 1.0))
        with f1:
            sel_ports = st.multiselect(
                "Portfolio (local filter for RCA visuals)",
                options=ports_options,
                default=ports_options,
                key="compl_rca_row2_ports",
            )
        with f2:
            sel_rca1 = st.multiselect(
                "Complaint Reason — RCA1 (local filter)",
                options=rca1_options,
                default=rca1_options,
                key="compl_rca_row2_rca1",
            )

        comp_local = comp.copy()
        if port_k and sel_ports:
            comp_local = comp_local[comp_local[port_k].astype(str).isin([str(p) for p in sel_ports])]

        rca2_filtered, rca1_filtered = _rca_tables_for_june(comp_local, use_ai=use_ai)
        if not rca1_filtered.empty and sel_rca1:
            rca1_filtered = rca1_filtered[rca1_filtered["RCA1"].isin(sel_rca1)]

        c3, c4 = st.columns((1.05, 1.0), gap="large")
        with c3:
            if not rca1_filtered.empty:
                _plot_rca1_pareto(rca1_filtered)
            else:
                _header("RCA1 — June 2025 (Pareto)")
                st.info("No data for the selected filters.")

        with c4:
            _header("RCA2 — June 2025 (by Portfolio)")
            rca2_by_port = _rca2_table_by_portfolio_for_june(
                comp=comp,
                use_ai=use_ai,
                portfolios=sel_ports if sel_ports else ports_options,
                rca1_keep=sel_rca1 if sel_rca1 else rca1_options,
            )
            if rca2_by_port.empty:
                st.info("No June-2025 complaints for the selected filters.")
            else:
                st.dataframe(_style_table(rca2_by_port[["Portfolio", "RCA2", "count"]], formats={"count": "{:,.0f}"}), use_container_width=True)

        rca2_all, rca1_all = _rca_tables_for_june(comp, use_ai=use_ai)
        if _PPT_READY and not table.empty and not mom.empty and not rca1_all.empty and not rca2_all.empty:
            ppt_bytes = _build_ppt(table, mom, rca1_all, rca2_all)
            st.download_button(
                "Download PPT",
                data=ppt_bytes,
                file_name="Complaint_Analysis_Jun2025.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
            )
        elif not _PPT_READY:
            st.caption("Install `python-pptx` to enable PPT download.")

    # --------- PORTFOLIO TABS ---------
    for i, portfolio in enumerate(portfolios, start=2):
        with tabs[i]:
            st.markdown(f"<h2 style='color:{_DARK_BLUE};margin:.3rem 0 1rem 0;'>{portfolio} — complaints analysis</h2>", unsafe_allow_html=True)

            t1, t2 = st.columns((1.1, 1.0), gap="large")
            with t1:
                df_trend = _portfolio_mom_series(cases, comp, portfolio)
                if not df_trend.empty:
                    st.pyplot(_fig_portfolio_trend(df_trend, f"{portfolio} trend Jan–Aug ’25"))
            with t2:
                df_reason_trend = _reason_trend_df(comp, portfolio, use_ai=use_ai)
                if not df_reason_trend.empty:
                    st.pyplot(_fig_reason_trend(df_reason_trend))

            b1, b2 = st.columns((1.1, 1.0), gap="large")
            with b1:
                df_delay = _delay_split_df(comp, portfolio, use_ai=use_ai)
                if not df_delay.empty:
                    st.pyplot(_fig_delay_split(df_delay))

            with b2:
                left, right = st.columns(2)
                with left:
                    _header("Delay — June (Top 80%)")
                    tdelay = _rca_tables_for_june(comp, use_ai=use_ai)[0]
                    if not tdelay.empty:
                        tdelay = tdelay.rename(columns={"RCA2":"Delay 80% Reason","percent":"Percentage contribution"})[["Delay 80% Reason","Percentage contribution"]]
                        st.dataframe(_style_table(tdelay, formats={"Percentage contribution": "{:.1f}%"}), use_container_width=True)
                with right:
                    _header("Procedure — June (Top 80%)")
                    # Recompute with filter = Procedure
                    _, port_k, date_k, desc_col = _detect_complaints_fields(comp)
                    df = comp.copy()
                    df["_month"] = _build_month_column(df, date_k, assume_year=2025 if date_k and date_k.lower() == "month" else None)
                    sub = df.loc[(df[port_k] == portfolio) & (df["_month"] == "2025-06"), [desc_col]].fillna("")
                    if not sub.empty:
                        texts = sub[desc_col].astype(str).tolist()
                        r1 = [_rca1_keyword(t) for t in texts]; r2 = [_rca2_keyword(t) for t in texts]
                        proc = [r2[i] for i in range(len(r2)) if r1[i] == "Procedure"]
                        if proc:
                            s = pd.Series(proc).value_counts().rename_axis("Procedure").reset_index(name="count")
                            total = max(1, s["count"].sum())
                            s["Percentage contribution"] = (s["count"]*100/total).round(0)
                            s = s[["Procedure","Percentage contribution"]]
                            st.dataframe(_style_table(s, formats={"Percentage contribution": "{:.0f}%"}), use_container_width=True)

    return ("", pd.DataFrame())


if __name__ == "__main__":
    print("Complaints module imports OK")
