from __future__ import annotations
import importlib
import traceback
from typing import Any, Dict, Optional, List, Tuple
from datetime import datetime
import io
import base64
import requests

import pandas as pd
import streamlit as st

# PPTX helpers
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# =============== Import helpers ===============
def _imp(mod: str, attr: str | None = None):
    """Import from repo root; if that fails, from core.<mod>."""
    try:
        m = importlib.import_module(mod)
    except ModuleNotFoundError:
        m = importlib.import_module(f"core.{mod}")
    return getattr(m, attr) if attr else m


# Your existing modules
load_store = _imp("data_store", "load_store")
sem_router = _imp("semantic_router")  # must define match(q) -> {"slug": ..., "params": {...}}

# Question modules are always looked up here
QUESTION_MODULE_PREFIXES = ("questions", "core.questions")


def _run_question(store: Dict[str, Any], slug: str, params: Dict[str, Any], user_text: Optional[str] = None):
    """Import and run a question module."""
    last_exc = None
    for prefix in QUESTION_MODULE_PREFIXES:
        try:
            mod = importlib.import_module(f"{prefix}.{slug}")
            return mod.run(store, params, user_text=user_text)
        except Exception as e:
            last_exc = e
            continue
    err = f"That question module failed to import.\n\nslug={slug}\n\n{traceback.format_exc()}"
    return err, pd.DataFrame()


# ===================== Page setup ====================
st.set_page_config(page_title="Halo - Quality - AI Assistant", layout="wide")

# ---------- Header styles ----------
st.markdown(
    """
    <style>
      [data-testid="stToolbar"] { display:none !important; }
      .halo-wrap{ display:flex; align-items:baseline; gap:.75rem; margin:8px 0 22px 0; }
      .halo-pill{
        background: linear-gradient(90deg, #FF7A00 0%, #FFD54F 50%, #66BB6A 100%);
        color:white; font-weight:900; letter-spacing:1.2px;
        border-radius:12px; padding:8px 16px; font-size:28px;
        text-shadow: 0 1px 1px rgba(0,0,0,.18);
      }
      .brand-title{ color:#0B3B8C; font-weight:800; font-size:36px; line-height:1.1; }
    </style>
    """,
    unsafe_allow_html=True,
)

# ---------- Header ----------
st.markdown(
    """
    <div class="halo-wrap">
      <div class="halo-pill">HALO</div>
      <div class="brand-title">- Quality - AI Assistant</div>
    </div>
    """,
    unsafe_allow_html=True,
)

# ============================== Load “store” ==============================
with st.spinner("Loading data..."):
    try:
        store = load_store(assume_year_for_complaints=2025)
    except TypeError:
        store = load_store()

# ============================== Router + query box ==============================
if "q" not in st.session_state:
    st.session_state["q"] = "fpa"   # default only once

q = st.text_input(
    "Type your question (e.g., 'comp', 'complaint', 'nps' or 'fpa')",
    key="q",
)

user_query = (q or "").strip()
match = sem_router.match(user_query) if hasattr(sem_router, "match") else {"slug": "complaints_june_by_portfolio", "params": {}}
slug = match.get("slug", "complaints_june_by_portfolio")
params = match.get("params", {}) or {}

# ============================== Conditional sidebar visibility ===============================
def _wants_sidebar(text: str, p: Dict[str, Any]) -> bool:
    if p.get("show_sidebar") is True:
        return True
    t = (text or "").lower()
    keywords = ("filter", "filters", "filter pane", "with filters", "show filters")
    return any(k in t for k in keywords)

SHOW_SIDEBAR = _wants_sidebar(user_query, params)
if not SHOW_SIDEBAR:
    st.markdown(
        """
        <style>
          [data-testid="stSidebar"], [data-testid="stSidebarNav"] { display: none !important; }
          section[data-testid="stSidebar"] { display: none !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )


# ============================== PPTX snapshot builder ===============================
def _load_logo_bytes() -> bytes | None:
    """Optional logo (base64 PNG) under [branding].logo_b64 in secrets."""
    b64 = st.secrets.get("branding", {}).get("logo_b64")
    if not b64:
        return None
    try:
        return base64.b64decode(b64)
    except Exception:
        return None


def _ppt_add_title(slide, title: str, subtitle: str | None, logo_bytes: bytes | None = None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.9)).text_frame
    tb.text = title
    p = tb.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(11, 61, 145)

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(8.5), Inches(0.5)).text_frame
    stamp = datetime.now().strftime("%d %b %Y %H:%M")
    sb.text = (subtitle or "Auto summary") + f"  •  {stamp}"
    sp = sb.paragraphs[0]
    sp.font.size = Pt(12)
    sp.font.color.rgb = RGBColor(90, 90, 90)

    if logo_bytes:
        bio = io.BytesIO(logo_bytes)
        bio.seek(0)
        slide.shapes.add_picture(bio, Inches(9.2), Inches(0.25), height=Inches(0.7))


def _add_caption(slide, text: str, left_in, top_in, width_in):
    cap = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.3)).text_frame
    cap.text = text
    cap.paragraphs[0].font.size = Pt(11)
    cap.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)


def _ppt_add_fig(slide, fig, left_in, top_in, width_in, caption=""):
    try:
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=220, bbox_inches="tight")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(left_in), Inches(top_in), width=Inches(width_in))
        if caption:
            _add_caption(slide, caption, left_in, top_in + 2.85, width_in)
        return top_in + 3.15
    except Exception:
        # silently skip bad figures
        return top_in


def _df_to_table(slide, df: pd.DataFrame, left_in, top_in, width_in, max_rows=8, caption=""):
    if not isinstance(df, pd.DataFrame) or df.empty:
        return top_in

    df = df.head(max_rows)
    rows, cols = len(df.index) + 1, len(df.columns)

    shape = slide.shapes.add_table(rows, cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(1.0))
    table = shape.table

    # Header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(11, 61, 145)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(240, 245, 255)

    # Body
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = "" if pd.isna(val) else str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
        # zebra
        if i % 2 == 0:
            for j in range(cols):
                c = table.cell(i, j)
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(250, 250, 250)

    # Even column widths
    for j in range(cols):
        table.columns[j].width = Inches(width_in / cols)

    if caption:
        _add_caption(slide, caption, left_in, top_in - 0.28, width_in)

    return top_in + 1.25


def _build_one_pager(
    title: str,
    subtitle: str | None,
    figs: List[Tuple[str, Any]] | None = None,
    tables: List[Tuple[str, pd.DataFrame]] | None = None,
    two_column: bool = True,
) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _ppt_add_title(slide, title, subtitle, logo_bytes=_load_logo_bytes())

    # Layout grid
    left_col, right_col = 0.5, 5.2
    col_w = 4.4
    y_left = y_right = 1.6

    def place_block(which: str, caption: str, payload):
        nonlocal y_left, y_right
        left = (y_left <= y_right) or not two_column
        x = left_col if left or not two_column else right_col
        y = y_left if left or not two_column else y_right

        if which == "fig":
            new_y = _ppt_add_fig(slide, payload, x, y, col_w, caption=caption)
        else:
            new_y = _df_to_table(slide, payload, x, y + 0.25, col_w, caption=caption)

        if two_column:
            if left:
                y_left = new_y + 0.15
            else:
                y_right = new_y + 0.15
        else:
            y_left = new_y + 0.15

    # Safely place figs/tables
    for cap, fig in (figs or []):
        if hasattr(fig, "savefig"):  # basic sanity check
            place_block("fig", cap or "", fig)

    for cap, df in (tables or []):
        if isinstance(df, pd.DataFrame):
            place_block("tbl", cap or "", df)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()


# ============================== Snapshot content helper ===============================
def _coerce_pairs(obj, expect_df: bool = False) -> List[Tuple[str, Any]]:
    """
    Ensures we always return a list of (caption, payload) tuples.
    - Accepts list of tuples, list of payloads, dict {caption: payload}, or None.
    - Filters invalid entries.
    """
    out: List[Tuple[str, Any]] = []
    if obj is None:
        return out

    if isinstance(obj, dict):
        it = list(obj.items())
    elif isinstance(obj, (list, tuple)):
        it = list(obj)
    else:
        # single payload -> wrap with default caption
        it = [("Item", obj)]

    for item in it:
        if isinstance(item, tuple) and len(item) == 2:
            cap, payload = item
        else:
            cap, payload = ("Item", item)

        if expect_df:
            if isinstance(payload, pd.DataFrame):
                out.append((str(cap), payload))
        else:
            # fig: we can't validate fully; check for savefig attribute
            if hasattr(payload, "savefig"):
                out.append((str(cap), payload))

    return out


def _get_snapshot_content(slug: str, store, params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Ask each question module for snapshot-ready content.
    Returns dict: {title, subtitle, figs, tables}
    with figs/tables **always** lists of (caption, payload) after coercion.
    """
    for prefix in QUESTION_MODULE_PREFIXES:
        try:
            mod = importlib.import_module(f"{prefix}.{slug}")
            # Preferred: build_snapshot(store, params)
            if hasattr(mod, "build_snapshot"):
                raw = mod.build_snapshot(store, params) or {}
                title = raw.get("title") or f"Halo Quality — {slug.replace('_', ' ').title()}"
                subtitle = raw.get("subtitle") or ""
                figs = _coerce_pairs(raw.get("figs"), expect_df=False)
                tables = _coerce_pairs(raw.get("tables"), expect_df=True)
                return {"title": title, "subtitle": subtitle, "figs": figs, "tables": tables}

            # Fallback: run() and use its DataFrame
            if hasattr(mod, "run"):
                result, df = mod.run(store, params)
                title = (result[0] if isinstance(result, tuple) and result else
                         result if isinstance(result, str) else
                         f"Halo Quality — {slug.replace('_', ' ').title()}")
                subtitle = (result[1] if isinstance(result, tuple) and len(result) > 1 else "")
                tables = _coerce_pairs([("Data", df)] if isinstance(df, pd.DataFrame) and not df.empty else [], expect_df=True)
                return {"title": title, "subtitle": subtitle, "figs": [], "tables": tables}
        except ModuleNotFoundError:
            continue
        except Exception as e:
            st.error(f"Snapshot content error in {slug}: {e}")

    # Final safe default
    return {"title": f"Halo Quality — {slug.replace('_', ' ').title()}", "subtitle": "Auto summary", "figs": [], "tables": []}


# ============================== Graph-based mail sender ===============================
def _get_graph_token(tenant_id: str, client_id: str, client_secret: str) -> str:
    token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    payload = {
        "client_id": client_id,
        "scope": "https://graph.microsoft.com/.default",
        "client_secret": client_secret,
        "grant_type": "client_credentials",
    }
    resp = requests.post(token_url, data=payload, timeout=10)
    if resp.status_code != 200:
        raise RuntimeError(f"Token request failed: {resp.status_code} {resp.text}")
    return resp.json().get("access_token")


def _send_via_graph(to_addrs: List[str], subject: str, html_body: str,
                    attachment_name: str, attachment_bytes: bytes):
    cfg = st.secrets.get("graph", {})
    tenant_id, client_id, client_secret, mailbox = (
        cfg.get("tenant_id"),
        cfg.get("client_id"),
        cfg.get("client_secret"),
        cfg.get("mailbox"),
    )
    if not all([tenant_id, client_id, client_secret, mailbox]):
        raise RuntimeError("Missing Graph configuration in secrets.toml under [graph].")

    token = _get_graph_token(tenant_id, client_id, client_secret)
    url = f"https://graph.microsoft.com/v1.0/users/{mailbox}/sendMail"

    attachment_b64 = base64.b64encode(attachment_bytes).decode("ascii")
    attachments = [
        {
            "@odata.type": "#microsoft.graph.fileAttachment",
            "name": attachment_name,
            "contentType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "contentBytes": attachment_b64,
        }
    ]

    message = {
        "message": {
            "subject": subject,
            "body": {"contentType": "HTML", "content": html_body},
            "toRecipients": [{"emailAddress": {"address": a}} for a in to_addrs],
            "attachments": attachments,
        },
        "saveToSentItems": "true",
    }

    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    resp = requests.post(url, json=message, headers=headers, timeout=15)
    if resp.status_code not in (200, 202):
        raise RuntimeError(f"Graph sendMail failed: {resp.status_code} {resp.text}")
    return True


# ============================== Share / Email snapshot UI ===============================
with st.expander("✉️  Share / Email snapshot"):
    default_recipients = st.secrets.get("graph", {}).get("recipients", []) or st.secrets.get("email", {}).get("recipients", [])
    selected = st.multiselect("Choose recipient(s)", options=default_recipients, default=default_recipients[:1] if default_recipients else [])
    custom = st.text_input("Or add another recipient (optional)", placeholder="name@company.com")

    col_a, col_b = st.columns([1, 1])
    with col_a:
        fmt = st.selectbox("Format", ["PPTX (one page)"], index=0)
    with col_b:
        include_tables = st.checkbox("Include small summary tables", value=True)

    if st.button("Send snapshot"):
        to_list = list(selected)
        if custom and "@" in custom:
            to_list.append(custom)

        if not to_list:
            st.error("Pick or enter at least one email address.")
        else:
            try:
                snap = _get_snapshot_content(slug, store, params)
                title = snap.get("title") or f"Halo Quality — {slug.replace('_', ' ').title()}"
                subtitle = snap.get("subtitle") or ""
                figs = snap.get("figs") or []
                tables = snap.get("tables") if include_tables else []

                # Build the PPT
                ppt_bytes = _build_one_pager(title, subtitle, figs=figs, tables=tables)

                # Send
                _send_via_graph(
                    to_addrs=to_list,
                    subject=title,
                    html_body=f"<p>Hi,</p><p>Attached is the snapshot from <b>Halo Quality</b> ({slug}).</p>",
                    attachment_name=f"{slug}_snapshot.pptx",
                    attachment_bytes=ppt_bytes,
                )
                st.success(f"Snapshot sent to: {', '.join(to_list)} ✅")
            except Exception as e:
                st.error(f"Could not send email: {e}")


# ============================== Run question ===============================
try:
    result, df = _run_question(store, slug, params, user_text=user_query)
except Exception:
    st.error("Sorry—couldn't run that question.")
    st.code(traceback.format_exc())
else:
    if isinstance(result, tuple) and len(result) in (1, 2):
        title = result[0]
        subtitle = result[1] if len(result) == 2 else None
        if isinstance(title, str) and title.strip():
            st.subheader(title)
        if subtitle:
            st.caption(subtitle)
    elif isinstance(result, str) and result.strip():
        st.info(result)

    if isinstance(df, pd.DataFrame) and not df.empty:
        st.dataframe(df, use_container_width=True)
